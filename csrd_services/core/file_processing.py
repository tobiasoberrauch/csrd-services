from csrd_services.core.esrs_classification import classify_text
from docling.document_converter import DocumentConverter
from langchain.text_splitter import MarkdownHeaderTextSplitter
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tempfile import NamedTemporaryFile


async def process_file(file, file_extension, contents):
    """
    Process the uploaded file based on its extension and return classification results.
    """
    with NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as tmp:
        tmp.write(contents)
        tmp_path = tmp.name

    converter = DocumentConverter()
    result = converter.convert(tmp_path)
    docs = result.document.export_to_markdown()

    md_splitter = MarkdownHeaderTextSplitter(
        headers_to_split_on=[("#", "Header_1"), ("##", "Header_2"), ("###", "Header_3")]
    )
    splits = md_splitter.split_text(docs)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500, separators=["\n\n", "\n"], chunk_overlap=0
    )
    splits = splitter.split_documents(splits)
    texts = [s.strip() for s in (i.page_content for i in splits)]

    return classify_text(texts)


async def process_pdf(contents):
    from langchain_community.document_loaders import PyPDFLoader
    import tempfile

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(contents)
        tmp_path = tmp.name

    loader = PyPDFLoader(tmp_path)
    docs = loader.load()
    text_content = [doc.page_content for doc in docs]
    return classify_text(text_content)


async def process_docx(contents):
    from docx import Document
    import tempfile

    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(contents)
        tmp_path = tmp.name

    doc = Document(tmp_path)
    text_content = [p.text for p in doc.paragraphs if p.text.strip()]
    return classify_text(text_content)


async def process_pptx(contents):
    from pptx import Presentation
    import tempfile

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(contents)
        tmp_path = tmp.name

    presentation = Presentation(tmp_path)
    text_content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content.append(shape.text)
    return classify_text(text_content)


async def process_xlsx(contents):
    import pandas as pd
    import tempfile

    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
        tmp.write(contents)
        tmp_path = tmp.name

    excel = pd.ExcelFile(tmp_path)
    text_content = {
        sheet: excel.parse(sheet).to_string(index=False) for sheet in excel.sheet_names
    }
    return classify_text(list(text_content.values()))


async def process_txt(contents):
    text_content = contents.decode("utf-8").splitlines()
    return classify_text(text_content)
